import tkinter as tk
from tkinter import *
from tkinter import filedialog
import customtkinter
import csv
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION
from datetime import date, timedelta, datetime
import os

other_care_type = [] 
perm_list = []
respite_list = []
initial_list = []
capped_list = []
uncapped_list = []
other_list = []
combined_cap_uncap = []
stay_type = [] 
perc_perm_resp = []
less_than_12m = []
more_than_12m = [] 
less_and_more_than_12m = []
combined_cap_uncap_qty = []
less_and_more_than_12m_qty = []
appraisal_type = []
adl_h = []
adl_m = []
adl_l = []
adl_n = []
beh_h = []
beh_m = []
beh_l = []
beh_n = []
chc_h = []
chc_m = []
chc_l = []
chc_n = []
adl_dom = []
beh_dom = []
chc_dom = []
perc_adl_dom = []
perc_beh_dom = []
perc_chc_dom = []

many = 0
count = 0
new_count = 0

customtkinter.set_default_color_theme("blue") 

#----Parent app config-----#
root = customtkinter.CTk()
root.title("Classy Claim Statements")
root.geometry("540x190") #WxH
root.resizable(False, False)


# ---Frame for widgets---#
mainframe = customtkinter.CTkFrame(root, height=200, width=600)
mainframe.grid(column=0, row=0, sticky=(N, S, E, W))
mainframe.columnconfigure(0, weight=1)
mainframe.rowconfigure(0, weight=1)


# Opens CSV claim statement.
def open_claim(): 
    global filename
    filename = filedialog.askopenfilename(
        initialdir="C:/Users/MainFrame/Desktop/", 
        title="Open CSV file", 
        filetypes=(("CSV Files", "*.csv"),)
    )
    pathh.insert(END, filename)
    print('file name', filename)

# Where presenation saves to. 
def export_place(): 
    global export_name
    export_name = filedialog.asksaveasfilename(
        initialdir="C:/Users/MainFrame/Desktop/", 
        title="Save PowerPoint Presentation To", 
        filetypes=(("PowerPoint Presentation", "*.pptx"),)
    )
    pathh2.insert(END, export_name)
    print('export name', export_name)

# Pull out data from csv(a long version)
def get_claim_data(): 
    global count
    global many
    global new_count

    try:
        one_year_ago = date.today() - timedelta(days = 365)
        with open(filename) as f:
            data = csv.reader(f)
            next(data, None)  # skip the headers
            for row in data:
            #--ACFI DOMAIN FILTER---
                count = count + 1
                if row[6] == 'DR':
                    new_count = new_count + 1
                try: #ADL domain
                    if 'H' in row[6][0]:
                        adl_h.append(row)
                    elif 'M' in row[6][0]:
                        adl_m.append(row)
                    elif 'L' in row[6][0]:
                        adl_l.append(row)
                    elif 'N' in row[6][0]:
                        adl_n.append(row)    
                    elif 'DR' in row[6]:#ignore new subs less than 1 month old
                        continue
                except IndexError: #respite
                    continue
                try: #Behaviour domain
                    if 'H' in row[6][1]:
                        beh_h.append(row)
                    elif 'M' in row[6][1]:
                        beh_m.append(row)
                    elif 'L' in row[6][1]:
                        beh_l.append(row)
                    elif 'N' in row[6][1]:
                        beh_n.append(row)
                    elif 'DR' in row[6]: #ignore new subs less than 1 month old 
                        continue
                except IndexError: #respite
                    continue  
                try:#Complex domain
                    if 'H' in row[6][2]:
                        chc_h.append(row)
                    elif 'M' in row[6][2]:
                        chc_m.append(row)
                    elif 'L' in row[6][2]:
                        chc_l.append(row)
                    elif 'N' in row[6][2]:
                        chc_n.append(row)
                    elif 'DR' in row[6]: #ignore new subs less than 1 month old
                        continue
                except IndexError: #respite
                    continue

                if row[4] == 'Respite': # respite people
                    respite_list.append(row)
                elif row[6] == 'DR' and row[9] != 'ACFI': # New, not subbed.
                    initial_list.append(row)
                elif row[6] == 'DR' and row[9] == 'ACFI': # New, subbed.
                    if row[10] == 'HHH': # New, capped
                        capped_list.append(row)
                    elif row[10] != 'HHH': # New, un-capped
                        uncapped_list.append(row)
                    else:
                        other_list.append(row) # abnormal
                elif row[6] != 'HHH' and row[6] != 'DR': # uncapped group
                    uncapped_list.append(row)
                elif row[6] == 'HHH': # capped group
                    capped_list.append(row)
                else: 
                    other_list.append(row) # abnormal group
                if row[4] == 'Permanent':
                    perm_list.append(row)
                if row[4] != 'Respite' and row[4] != 'Permanent':
                    other_care_type.append(row) #incase there is something odd

            # refining the data/calculations
            appraisal_type.append(len(perm_list) - new_count)
            appraisal_type.append(len(respite_list))
            appraisal_type.append(new_count)
            appraisal_type.append(len(other_care_type))
            all_claims = len(uncapped_list) + len(capped_list)
            perc_uncap = str(round((len(uncapped_list) / all_claims) * 100))
            perc_cap = str(round((len(capped_list) / all_claims) * 100))
            combined_cap_uncap.append(round(int(perc_cap))) # capped number first
            combined_cap_uncap.append(round(int(perc_uncap))) # uncapped second
            combined_cap_uncap_qty.append(len(capped_list)) # capped number first
            combined_cap_uncap_qty.append(len(uncapped_list)) # uncapped second
            
        # sort uncapped residents    
        for dates in uncapped_list:
            if dates[6] != 'DR': # We filter DR's as they are always < 12momnths & cause error.
                if '/' not in dates[7]: # if no date, then it's a mandatory which is always < 12 months.
                    try:
                        dates_list = datetime.strptime(dates[8], '%d/%m/%Y').date() #convert string date to datetime object
                        if dates_list < one_year_ago: 
                            many += 1
                            name = str(dates[1]) + ', ' + str(dates[2]) + ', ' + str(dates[6]) + ', Exp date: ' + str(dates[8]) 
                            more_than_12m.append(name) 
                        else:
                            less_than_12m.append(dates)
                    except ValueError: less_than_12m.append(dates)
                else:
                    less_than_12m.append(dates) # Mandatory < 12 months.
            else:
                less_than_12m.append(dates) # DR's are always < 12 months.

        a1 = len(uncapped_list) 
        a2 = len(less_than_12m) 
        perc_old = a2 / a1 * 100 #precent of uncapped more than 12 months old
        b1 = len(uncapped_list) 
        b2 = len(more_than_12m) 
        perc_new = b2 / b1 * 100 #precent of uncapped less than 12 months old
        less_and_more_than_12m.append(round(perc_old))
        less_and_more_than_12m.append(round(perc_new))
        less_and_more_than_12m_qty.append(len(less_than_12m))
        less_and_more_than_12m_qty.append(len(more_than_12m))
        perm = count - len(respite_list)
        perc_perm_resp.append(round((perm/count)*100))
        perc_perm_resp.append(round((len(respite_list)/count)*100))
        stay_type.append(perm)
        stay_type.append(len(respite_list))
        adl_dom.append(len(adl_h))
        adl_dom.append(len(adl_m))
        adl_dom.append(len(adl_l))
        adl_dom.append(len(adl_n))
        beh_dom.append(len(beh_h))
        beh_dom.append(len(beh_m))
        beh_dom.append(len(beh_l))
        beh_dom.append(len(beh_n))
        chc_dom.append(len(chc_h))
        chc_dom.append(len(chc_m))
        chc_dom.append(len(chc_l))
        chc_dom.append(len(chc_n))

        #get percentage of each domain category.
        perc_adl_dom.append(round((len(adl_h)/(count - ((len(respite_list)) + new_count)))*100)) 
        perc_adl_dom.append(round((len(adl_m)/(count - ((len(respite_list)) + new_count)))*100))
        perc_adl_dom.append(round((len(adl_l)/(count - ((len(respite_list)) + new_count)))*100))
        perc_adl_dom.append(round((len(adl_n)/(count - ((len(respite_list)) + new_count)))*100))
        perc_beh_dom.append(round((len(beh_h)/(count - ((len(respite_list)) + new_count)))*100))
        perc_beh_dom.append(round((len(beh_m)/(count - ((len(respite_list)) + new_count)))*100))
        perc_beh_dom.append(round((len(beh_l)/(count - ((len(respite_list)) + new_count)))*100))
        perc_beh_dom.append(round((len(beh_n)/(count - ((len(respite_list)) + new_count)))*100))
        perc_chc_dom.append(round((len(chc_h)/(count - ((len(respite_list)) + new_count)))*100))
        perc_chc_dom.append(round((len(chc_m)/(count - ((len(respite_list)) + new_count)))*100))
        perc_chc_dom.append(round((len(chc_l)/(count - ((len(respite_list)) + new_count)))*100))
        perc_chc_dom.append(round((len(chc_n)/(count - ((len(respite_list)) + new_count)))*100))

        # call function to create presentaiton slide.
        create_presentation()

    # catch error if file not selected
    except NameError:
        tk.messagebox.showinfo('Oops', 'Please make sure both the file and location have been selected')
        print('you need to select file')        
        
#construct the presentation
def create_presentation():
    global prs 
    prs = Presentation()
    # order slides are in PowerPoint
    title_slide()
    appraisal_type_slide()
    cap_vs_uncap_slide()
    cap_vs_uncap_perc_slide()
    less_more_12m_qty_slide()
    less__more_12months_perc_slide()
    perm_vs_resp_qty_slide()
    perm_resp_percent_slide()
    adl_slide()
    adl_perc_slide()
    beh_slide()
    beh_perc_slide()
    chc_slide()
    chc_perc_slide()    
    closing_slide()
    

    export_to = export_name + '.pptx' 
    prs.save(export_to)

    # prompt and open created presentation
    tk.messagebox.showinfo('Presentation Created', 'Click OK to close this window and open your beautiful presentation')
    os.startfile(export_to)

    # close GUI window
    root.destroy()
    return 
    
#title slide
def title_slide(): 
    slide17 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide17.shapes.title
    subtitle = slide17.placeholders[1]
    title.text = 'Medicare Claim Statement' + '\n' + 'Data Visualisation.'
    subtitle.text = 'Ready for you to style.'
    return

#appraisal type slide
def appraisal_type_slide(): 
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('Permanent', 'Respite', 'New Resident(< 60 days)', 'Other')
    chart_data.add_series('Appraisal Types', appraisal_type )
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart12 = slide3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart12.plots[0].has_data_labels = True
    data_labels = chart12.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    return

# QTY resident capped and uncapped
def cap_vs_uncap_slide(): 
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('Capped', 'Uncapped')
    chart_data.add_series('Capped Vs Uncapped Residents', combined_cap_uncap_qty )
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart2 = slide2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart2.plots[0].has_data_labels = True
    data_labels = chart2.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    return

# % resident capped and uncapped
def cap_vs_uncap_perc_slide(): 
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('Capped', 'Uncapped')
    chart_data.add_series(' % Capped Vs Uncapped Appraisals', combined_cap_uncap)
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart = slide6.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = '0\%'
    return

# QTY residents less and more than 12 months old.
def less_more_12m_qty_slide():  
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = (' Less than 12 months old', 'More than 12 months old')
    chart_data.add_series(" Age of Uncapped Appraisals", less_and_more_than_12m_qty)
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart9 = slide9.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart9.plots[0].has_data_labels = True
    data_labels = chart9.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    return

# % residents less and more than 12 months old.
def less__more_12months_perc_slide():  
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = (' Less than 12 months old', 'More than 12 months old')
    chart_data.add_series(" % uncapped", less_and_more_than_12m)
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart = slide15.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = '0\%'
    return

# QTY perm Vs Respite
def perm_vs_resp_qty_slide(): 
    slide31 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('Perm', 'Respite')
    chart_data.add_series(' Permanent Vs Respite', stay_type)
    # add chart to slide --------------------
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart1 = slide31.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart1.plots[0].has_data_labels = True
    data_labels = chart1.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    return

# % perm Vs Respite
def perm_resp_percent_slide(): 
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('Permanent', 'Respite')
    chart_data.add_series(" % Permanant Vs Respite", perc_perm_resp)
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart = slide5.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = '0\%'
    return

# ADL breakdown
def adl_slide(): 
    slide33 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('High', 'Medium', 'Low', 'Nil')
    chart_data.add_series('ADL Domain', adl_dom)
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart12 = slide33.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart12.plots[0].has_data_labels = True
    data_labels = chart12.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    return

# % adl domain
def adl_perc_slide(): 
    slide56 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('High', 'Medium', 'Low', 'Nil')
    chart_data.add_series("% ADL Domain", perc_adl_dom)
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart = slide56.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = '0\%'
    return

# Beh domain
def beh_slide(): 
    slide44 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('High', 'Medium', 'Low', 'Nil')
    chart_data.add_series('Behaviour Domain', beh_dom)
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart12 = slide44.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart12.plots[0].has_data_labels = True
    data_labels = chart12.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    return

# % beh domain
def beh_perc_slide(): 
    slide561 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('High', 'Medium', 'Low', 'Nil')
    chart_data.add_series("% Behaviour Domain", perc_beh_dom)
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart = slide561.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = '0\%'
    return

# chc domain
def chc_slide(): 
    slide55 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('High', 'Medium', 'Low', 'Nil')
    chart_data.add_series('Complex Health Care Domain', chc_dom)
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart12 = slide55.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
    chart12.plots[0].has_data_labels = True
    data_labels = chart12.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    return

# % chc domain
def chc_perc_slide(): 
    slide566 = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ('High', 'Medium', 'Low', 'Nil')
    chart_data.add_series("% Complex Health Care Domain", perc_chc_dom)
    x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
    chart = slide566.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
    chart.has_legend = True
    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = '0\%'
    return

def closing_slide(): 
    slide16 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide16.shapes.title
    subtitle = slide16.placeholders[1]
    title.text = 'This is a list of people who did not fit through the filter i.e exceptions/NQR on Claim Statement. (Quality check. Delete slide)'
    subtitle.text = str(other_list)
    return 



#------------GUI---------------#
select_file_label = customtkinter.CTkLabel(mainframe, 
                                            text="Select claim statement:",
                                            width=100, 
                                            height=50)
select_file_label.grid(row=0, column=0, pady=10, padx=30, sticky=(W))

pathh = customtkinter.CTkEntry(mainframe, width=140)
pathh.grid(row=0, column=1, columnspan=2, pady=5, padx=10, sticky=())


select_file_btn1 = customtkinter.CTkButton(mainframe, 
                                            text="Select file", 
                                            command=open_claim)
select_file_btn1.grid(row=0, column=3, pady=5, padx=10)


# -------CHOOSE PRESENTATION SAVE LOCATION---#
export_label = customtkinter.CTkLabel(mainframe, 
                                            text="Save presentation to:",
                                            width=100, 
                                            height=50)
                                            #fg_color=None)
export_label.grid(row=1, column=0, pady=10, padx=30, sticky=(W))

pathh2 = customtkinter.CTkEntry(mainframe, width=140)
pathh2.grid(row=1, column=1, columnspan=2, pady=5, padx=10, sticky=())


export_btn1 = customtkinter.CTkButton(mainframe, 
                                            text="Select location", 
                                            command=export_place)
export_btn1.grid(row=1, column=3, pady=5, padx=10)



# --------CREATE PRESENTATION button-------
go_btn1 = customtkinter.CTkButton(mainframe, 
                                            text="Create Presentation", 
                                            command=get_claim_data)
go_btn1.grid(row=2, column=2, pady=5, padx=10)




# run GUI
root.mainloop()
